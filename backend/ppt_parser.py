import os

from pptx import Presentation
import fitz


def extract_from_ppt(path):

    presentation = Presentation(path)

    text = ""

    for slide in presentation.slides:

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text += shape.text + "\n"

    return text


def extract_from_pdf(path):

    doc = fitz.open(path)

    text = ""

    for page in doc:

        text += page.get_text()

    return text


def extract_text(path):

    ext = os.path.splitext(path)[1].lower()

    if ext == ".pptx":
        return extract_from_ppt(path)

    if ext == ".pdf":
        return extract_from_pdf(path)

    raise Exception(
        "Only PDF and PPTX supported"
    )